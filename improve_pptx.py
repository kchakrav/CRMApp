#!/usr/bin/env python3
"""
Improve PowerPoint (.pptx) presentation: fonts, colors, layout, and consistency.
Usage: python improve_pptx.py [input.pptx] [output.pptx]
Default: reads from user's Desktop, saves with _improved suffix.
"""

import sys
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Installing python-pptx...")
    os.system(f"{sys.executable} -m pip install python-pptx -q")
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Design constants
FONT_TITLE = "Segoe UI"
FONT_BODY = "Segoe UI"
SIZE_TITLE = Pt(32)
SIZE_BODY = Pt(18)
SIZE_SUBTITLE = Pt(14)
COLOR_TITLE = RgbColor(0x2C, 0x3E, 0x50)      # Dark slate
COLOR_BODY = RgbColor(0x34, 0x49, 0x5E)        # Slightly lighter
COLOR_ACCENT = RgbColor(0x34, 0x98, 0xDB)      # Blue accent
COLOR_SUBTITLE = RgbColor(0x7F, 0x8C, 0x8D)    # Gray


def safe_set_font(run, name, size, color):
    """Set font properties, skipping theme/unsupported colors."""
    try:
        run.font.name = name
        run.font.size = size
        run.font.color.rgb = color
    except (AttributeError, TypeError):
        pass


def improve_slide(slide):
    """Apply consistent formatting to all shapes on a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    safe_set_font(run, FONT_BODY, SIZE_BODY, COLOR_BODY)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            safe_set_font(run, FONT_BODY, SIZE_BODY, COLOR_BODY)


def improve_title_slide(slide):
    """Format title slide (usually first slide)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            for i, para in enumerate(tf.paragraphs):
                for run in para.runs:
                    run.font.name = FONT_TITLE
                    if i == 0:
                        run.font.size = Pt(44)
                        run.font.bold = True
                        safe_set_font(run, FONT_TITLE, Pt(44), COLOR_TITLE)
                    else:
                        run.font.size = SIZE_SUBTITLE
                        safe_set_font(run, FONT_TITLE, SIZE_SUBTITLE, COLOR_SUBTITLE)


def improve_presentation(prs):
    """Apply improvements across the entire presentation."""
    for i, slide in enumerate(prs.slides):
        if i == 0:
            improve_title_slide(slide)
        else:
            improve_slide(slide)
    return prs


def main():
    default_input = Path.home() / "OneDrive - Adobe" / "Desktop" / "Vibe coding with AI.pptx"
    if len(sys.argv) >= 2:
        input_path = Path(sys.argv[1])
    else:
        input_path = default_input

    if not input_path.exists():
        print(f"Input file not found: {input_path}")
        print("\nUsage: python improve_pptx.py [input.pptx] [output.pptx]")
        print("Example: python improve_pptx.py \"C:\\Users\\You\\Desktop\\Vibe coding with AI.pptx\"")
        sys.exit(1)

    if len(sys.argv) >= 3:
        output_path = Path(sys.argv[2])
    else:
        stem = input_path.stem
        output_path = input_path.parent / f"{stem}_improved.pptx"

    print(f"Reading: {input_path}")
    prs = Presentation(str(input_path))
    improve_presentation(prs)
    prs.save(str(output_path))
    print(f"Saved:  {output_path}")
    print("Done.")


if __name__ == "__main__":
    main()
